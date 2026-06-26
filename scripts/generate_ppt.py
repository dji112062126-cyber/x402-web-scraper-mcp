#!/usr/bin/env python3
"""Generate AI Report PPT for June 2026."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x11, 0x1A)   # near-black
BG_SECTION = RGBColor(0x12, 0x18, 0x24)   # dark card
ACCENT     = RGBColor(0x6C, 0x5C, 0xFF)   # purple accent
ACCENT2    = RGBColor(0x00, 0xBF, 0xA5)   # teal
ACCENT3    = RGBColor(0xFF, 0x6B, 0x6B)   # coral
ACCENT4    = RGBColor(0xFF, 0xD9, 0x3D)   # gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC4)
DIM_GRAY   = RGBColor(0x6B, 0x72, 0x80)
GRID_LINE  = RGBColor(0x1E, 0x26, 0x33)


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── helpers ────────────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height, lines):
    """lines: list of (text, font_size, color, bold) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text, font_size, color, bold = line_data
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=14, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(3)
    return txBox


def section_header(slide, title, subtitle=""):
    # accent bar at top
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=DIM_GRAY)


def card(slide, left, top, width, height, title, content_lines, accent_color=ACCENT):
    """A card with title + bullet content."""
    # card background
    add_rect(slide, left, top, width, height, fill_color=BG_SECTION)
    # accent strip on left
    add_rect(slide, left, top + Inches(0.05), Inches(0.06), height - Inches(0.1),
             fill_color=accent_color)
    # title
    add_text_box(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.35), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    # content
    add_bullet_points(slide, left + Inches(0.18), top + Inches(0.5),
                      width - Inches(0.35), height - Inches(0.65),
                      content_lines, font_size=12, color=LIGHT_GRAY)


def page_number(slide, num):
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                 str(num), font_size=11, color=DIM_GRAY, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)
# decorative shapes
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=ACCENT)
add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), fill_color=ACCENT3)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
             "2026 年 6 月", font_size=20, color=DIM_GRAY)

add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
             "AI 行业最新资讯报告", font_size=48, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "智能体时代全面到来 · 模型格局深度解析 · 行业风口成功案例", font_size=18, color=ACCENT2)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "来源: Computex 2026 / Microsoft Build / Snowflake Summit / Forbes AI 50 / Artificial Analysis",
             font_size=11, color=DIM_GRAY)

page_number(slide, 1)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — 目录
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "目  录", "CONTENTS")

toc_items = [
    ("01", "AI Agent 最新动态", "从 Computex 到 Microsoft Build，Agent 全面接管数字生活"),
    ("02", "Agent 框架全景", "LangChain · CrewAI · AutoGPT · AutoGen … 2026 主流框架对比与选型"),
    ("03", "Agent 协议标准化", "MCP vs A2A — AI Agent 生态的 TCP/IP 时刻"),
    ("04", "大模型强弱对比", "Opus 4.8 vs GPT-5.5 vs Gemini 3.1 全方位 Benchmark 对比"),
    ("05", "AI 风口行业 & 成功案例", "金融 · 医疗 · 教育 · 情感AI · 机器人 · 出海"),
    ("06", "行业趋势 & 未来展望", "2026 关键判断与战略建议"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(i * 0.92)
    add_text_box(slide, Inches(1.2), y, Inches(0.8), Inches(0.4),
                 num, font_size=28, color=ACCENT, bold=True)
    add_text_box(slide, Inches(2.1), y, Inches(9), Inches(0.38),
                 title, font_size=22, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.1), y + Inches(0.38), Inches(9), Inches(0.3),
                 desc, font_size=13, color=DIM_GRAY)

page_number(slide, 2)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — AI Agent 最新动态
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "01  AI Agent 最新动态", "2026 年 6 月 — 智能体时代全面到来")

# Key insight bar
add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7), fill_color=BG_SECTION)
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(11), Inches(0.6),
             '🔥  2026 被定义为 "Agent 元年" — NVIDIA CEO 黄仁勋: AI 已从生成式 AI 全面进入 Agentic AI 时代',
             font_size=15, color=ACCENT4, bold=True)

# Row 1 — 3 cards
card(slide, Inches(0.8), Inches(2.5), Inches(3.7), Inches(2.2),
     "🖥️ Computex 2026 台北",
     ["NVIDIA RTX Spark: Arm+Blackwell+128GB 统一内存", "Vera Rubin: 首个专为 Agent 设计的架构投产", "Nemotron 3 Ultra: 推理速度 5×↑ 成本降 30%", "高通: 2026 是 Agent 之年，Agent 成数字生活中心", "Intel Arc G3 + OpenVINO 物理 AI 开源机器人库"],
     accent_color=ACCENT)

card(slide, Inches(4.8), Inches(2.5), Inches(3.7), Inches(2.2),
     "☁️ Microsoft Build 2026",
     ["MAI-Thinking-1: 微软首个推理模型，企业级", "Copilot Super App 确认开发中，夏末预览", "Agent Mode 成为 M365 Copilot 默认模式", "Azure AI Foundry 深度支持多模型部署", "GitHub Copilot 覆盖完整开发工作流"],
     accent_color=ACCENT2)

card(slide, Inches(8.8), Inches(2.5), Inches(3.7), Inches(2.2),
     "🏢 更多巨头动态",
     ["Meta 发布 Business Agent，100 万+企业入驻", "Snowflake 推 Cortex Code + Intelligence 平台", "Anthropic Claude Code 支持数百子 Agent 动态编排", "Google 发布 Omni(any-to-any)，推 Antigravity 平台", "Cloudflare: Agentic 流量首次超过人类流量"],
     accent_color=ACCENT3)

# Bottom key data
add_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.8), fill_color=BG_SECTION)
add_text_box(slide, Inches(1.0), Inches(5.1), Inches(11), Inches(0.35),
             "📊 关键数据", font_size=16, color=WHITE, bold=True)
stats = [
    ["Agent 流量占比", "超 50%", "Agent 流量首次超过人类网页流量"],
    ["MCP 生态", "10,000+ 服务器 / 9700 万月下载", "A2A 合并 ACP 规范，100+公司支持"],
    ["AI Agent 框架", "400+ 预置工具", "LangChain 生态最大；CrewAI 多 Agent 协作领先"],
    ["企业 AI 支出", "3056 亿美元 (AI 50 合计融资)", "Anthropic+OpenAI 占 80%"],
]
for j, (label, value, note) in enumerate(stats):
    x = Inches(1.0) + Inches(j * 2.85)
    add_text_box(slide, x, Inches(5.45), Inches(2.7), Inches(0.25), label, font_size=11, color=DIM_GRAY)
    add_text_box(slide, x, Inches(5.68), Inches(2.7), Inches(0.3), value, font_size=18, color=ACCENT2, bold=True)
    add_text_box(slide, x, Inches(6.0), Inches(2.7), Inches(0.5), note, font_size=10, color=LIGHT_GRAY)

page_number(slide, 3)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Agent 框架全景
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "02  Agent 框架全景", "2026 主流框架对比 · 从链式到图式 · 多 Agent 协作成熟")

# Framework comparison table
header_data = [("框架", 2.2), ("核心优势", 3.3), ("最佳场景", 3.3), ("许可证", 1.5), ("参考价格", 1.5)]
x_start = Inches(0.8)
y_table = Inches(1.7)

add_rect(slide, x_start, y_table, Inches(11.7), Inches(0.45), fill_color=ACCENT)
for j, (hdr, w) in enumerate(header_data):
    x = x_start + sum(header_data[k][1] for k in range(j))
    add_text_box(slide, x + Inches(0.1), y_table, Inches(w - 0.2), Inches(0.45),
                 hdr, font_size=13, color=WHITE, bold=True)

rows = [
    ("LangChain / LangGraph", "最大生态 (400+ 工具)；图式工作流引擎", "RAG / 聊天机器人 / 企业工作流", "开源", "免费"),
    ("CrewAI", "多 Agent 角色协作；任务编排四阶段模型", "内容生产 / 多步骤业务流程", "开源 + 付费", "$99/月起"),
    ("Microsoft AutoGen", "对话驱动多 Agent；递归反思；并行执行", "研究报告 / 代码生成 / 复杂分析", "MIT 开源", "免费"),
    ("Semantic Kernel", "企业治理；Azure 集成；.NET+Python", "金融 / 医疗受监管行业", "开源", "免费 + Azure"),
    ("OpenAI Agents SDK", "原生 OpenAI 生态；工具调用最轻量", "GPT 原生应用 / 快速原型", "按 Token", "按量计费"),
    ("AutoGPT", "首个自主目标驱动 Agent；社区庞大", "自主实验 / 兴趣探索", "开源", "免费"),
    ("LlamaIndex / Mastra / Pydantic", "数据密集 RAG / TS 优先 / 强类型", "文档问答 / JS 团队 / Pydantic 用户", "开源", "免费"),
]

for i, row in enumerate(rows):
    y = y_table + Inches(0.5) + Inches(i * 0.62)
    bg_c = BG_SECTION if i % 2 == 0 else BG_DARK
    add_rect(slide, x_start, y, Inches(11.7), Inches(0.58), fill_color=bg_c)
    for j, (cell, w) in enumerate(zip(row, [w for _, w in header_data])):
        x = x_start + sum(header_data[k][1] for k in range(j))
        is_first = (j == 0)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.02), Inches(w - 0.2), Inches(0.52),
                     cell, font_size=11, color=WHITE if is_first else LIGHT_GRAY,
                     bold=is_first)

# Bottom insight
add_rect(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9), fill_color=BG_SECTION)
add_bullet_points(slide, Inches(1.0), Inches(6.1), Inches(11), Inches(0.7), [
    "趋势: 链式执行 → 动态执行图 (分支/循环/合并)；工具自动发现与自适应调用成为标配",
    "挑战: 复杂任务端到端成功率 ≤85%；成本控制与可靠性仍是瓶颈",
], font_size=12, color=LIGHT_GRAY)

page_number(slide, 4)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Agent 协议标准化 (MCP vs A2A)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "03  Agent 协议标准化", "MCP vs A2A — AI 世界的 TCP/IP 时刻")

# Two columns
card(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(3.2),
     "🔌 MCP (Model Context Protocol)",
     ["创建方: Anthropic (2024.11)，2025年捐赠给 Linux Foundation AAIF",
      "定位: AI 模型的「USB-C 接口」— 任何模型接入任何工具",
      "架构: Client → Server (JSON-RPC 2.0 over stdio/HTTP)",
      "生态: 10,000+ MCP 服务器，9700 万月 SDK 下载",
      "支持: Claude / ChatGPT / Gemini / VS Code / Cursor",
      "最新: 2025-06-18 版，新增 OAuth 2.0 资源指示器"],
     accent_color=ACCENT)

card(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(3.2),
     "🤝 A2A (Agent-to-Agent Protocol)",
     ["创建方: Google (2025.04)，2025年捐赠给 Linux Foundation",
      "定位: Agent 之间的「HTTP」— 跨厂商 Agent 发现与协作",
      "架构: Peer → Peer (HTTP + JSON-RPC + SSE; Agent Card 发现)",
      "生态: 100+ 公司支持 (Google/AWS/MS/Cisco/SAP/Salesforce)",
      "里程碑: v1.0 达成，合并 IBM ACP 协议",
      "特性: 任务状态机 (SUBMITTED→WORKING→COMPLETED/FAILED)"],
     accent_color=ACCENT2)

# Consensus
add_rect(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.7), fill_color=BG_SECTION)
add_text_box(slide, Inches(1.0), Inches(5.3), Inches(11), Inches(0.35),
             "✅ 行业共识: 互补而非竞争 — 两者并用", font_size=18, color=ACCENT4, bold=True)

consensus = [
    "MCP: 解决「我的 Agent 如何调用工具/数据？」 — 单 Agent 工具接入标准",
    "A2A: 解决「我的 Agent 如何与另一个 Agent 对话？」 — 多 Agent 跨厂商协作标准",
    "典型架构: 工具调用 → MCP，跨团队/跨厂商 Agent 协作 → A2A，统一网关做安全管控 (OAuth + Prompt Injection 防护)",
]
add_bullet_points(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(1.0), consensus, font_size=13, color=LIGHT_GRAY)

page_number(slide, 5)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — 大模型对比总览
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "04  大模型强弱对比 (一)", "Claude Opus 4.8 vs GPT-5.5 vs Gemini 3.1 Pro — 全方位 Benchmark 对比")

# Model comparison cards - 3 columns
card(slide, Inches(0.8), Inches(1.6), Inches(3.7), Inches(2.5),
     "🟠 Claude Opus 4.8 (Anthropic)",
     ["💻 SWE-bench Pro: 69.2% 🥇 (领先 GPT 10.6%)",
      "🧠 推理 (LiveBench): 89.71 🥇",
      "🤖 Agent (ARC-AGI-3): 1.5% 🥇 (4× 对手)",
      "📋 实时网页 (Online-Mind2Web): 84% 🥇",
      "🎯 可靠性: 幻觉率 35.9% 🥇 (最诚实)", "💰 $5 / $25 每百万 token"],
     accent_color=RGBColor(0xD4, 0x7B, 0x3A))

card(slide, Inches(4.8), Inches(1.6), Inches(3.7), Inches(2.5),
     "🟢 GPT-5.5 (OpenAI)",
     ["📐 数学 (LiveBench): 96.32 🥇",
      "💻 Terminal-Bench: 82.7% 🥇",
      "🧠 GPQA Diamond: ~93.6%",
      "🏆 LiveBench 全科: 80.71 🥇",
      "👁️ 多模态/视觉: 领先 Claude",
      "💰 $5 / $30 每百万 token"],
     accent_color=RGBColor(0x10, 0xA3, 0x7F))

card(slide, Inches(8.8), Inches(1.6), Inches(3.7), Inches(2.5),
     "🔵 Gemini 3.1 Pro (Google)",
     ["📋 指令遵循 (LiveBench IF): 79.10 🥇",
      "👁️ ScienceQA Vision: 61.86% 🥇",
      "🏆 LiveCodeBench Elo: 2887 🥇",
      "💻 SWE-bench Verified: 80.6% (Top 3)",
      "💰 性价比最优: $2 / $12 每百万 token",
      "📏 1M 上下文窗口"],
     accent_color=RGBColor(0x42, 0x85, 0xF4))

# Bottom comparison table
y_comp = Inches(4.5)
add_rect(slide, Inches(0.8), y_comp, Inches(11.5), Inches(2.4), fill_color=BG_SECTION)
add_text_box(slide, Inches(1.0), y_comp + Inches(0.05), Inches(11), Inches(0.3),
             "📊 核心维度雷达对比", font_size=15, color=WHITE, bold=True)

dims = [
    ("Agent/自主工作 🥇", "Claude"), ("编码/修Bug 🥇", "Claude"), ("数学 🥇", "GPT-5.5"),
    ("通用全科 🥇", "GPT-5.5"), ("指令遵循 🥇", "Gemini"), ("视觉/多模态 🥇", "Gemini"),
    ("性价比 🥇", "Gemini"), ("可靠性/诚实 🥇", "Claude"),
]
for j, (dim, winner) in enumerate(dims):
    x = Inches(0.9) + Inches(j % 4 * 2.95)
    y = y_comp + Inches(0.45) + Inches(j // 4 * 0.85)
    clr = ACCENT if winner == "Claude" else ACCENT2 if winner == "GPT-5.5" else ACCENT3
    add_text_box(slide, x, y, Inches(2.8), Inches(0.25), dim, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide, x, y + Inches(0.25), Inches(2.8), Inches(0.2), f"🏆 {winner}", font_size=11, color=clr, bold=True)

page_number(slide, 6)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — 大模型对比 (开源 & 性价比)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "04  大模型强弱对比 (二)", "开源模型 · 定价对比 · 选型建议")

# Pricing table
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.3),
             "💰 定价对比 (每百万 Token)", font_size=16, color=WHITE, bold=True)

price_data = [
    ("模型", "输入", "输出", "类型"),
    ("DeepSeek V4-Flash", "$0.04", "$0.07", "开源 / MIT"),
    ("Gemini 3.1 Flash", "$0.30", "$2.50", "商业"),
    ("DeepSeek V4-Pro", "$0.44", "$0.87", "开源 / MIT"),
    ("Gemini 3.1 Pro", "$2.00", "$12.00", "商业"),
    ("GPT-5.4", "$2.50", "$15.00", "商业"),
    ("Claude Sonnet 4.6", "$3.00", "$15.00", "商业"),
    ("Claude Opus 4.8", "$5.00", "$25.00", "商业"),
    ("GPT-5.5", "$5.00", "$30.00", "商业"),
]

for i, row in enumerate(price_data):
    y = Inches(1.85) + Inches(i * 0.42)
    is_header = (i == 0)
    bg = ACCENT if is_header else (BG_SECTION if i % 2 == 0 else BG_DARK)
    add_rect(slide, Inches(0.8), y, Inches(5.6), Inches(0.38), fill_color=bg)
    for j, cell in enumerate(row):
        x = Inches(0.9) + Inches(j * 1.3)
        fc = WHITE if is_header else LIGHT_GRAY
        add_text_box(slide, x, y + Inches(0.02), Inches(1.2), Inches(0.34),
                     cell, font_size=11, color=fc, bold=is_header)

# Open weight models
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.3),
             "🔓 开源/开放权重模型 Top 5", font_size=16, color=WHITE, bold=True)

open_models = [
    ("MiniMax M2.5", "SWE-bench 80.2% · $0.30/$1.20"),
    ("Kimi K2.6", "AA Intelligence 开源 #1 · ~$2 hosted"),
    ("DeepSeek V4-Pro", "SWE-bench 80.6% · MIT 许可"),
    ("Qwen 3.6-72B", "68.2% · Apache 2.0 · 可自部署"),
    ("GLM-5", "SWE-bench Pro 62.8% · MIT · 免费"),
]
for j, (name, desc) in enumerate(open_models):
    y = Inches(1.93) + Inches(j * 0.46)
    add_text_box(slide, Inches(7.2), y, Inches(2.5), Inches(0.22), name, font_size=13, color=ACCENT2, bold=True)
    add_text_box(slide, Inches(9.5), y, Inches(3.3), Inches(0.30), desc, font_size=10, color=LIGHT_GRAY)

# Best pick by use case
add_rect(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), fill_color=BG_SECTION)
add_text_box(slide, Inches(1.0), Inches(5.45), Inches(11), Inches(0.3),
             "🎯 按场景最佳选择", font_size=16, color=WHITE, bold=True)

picks = [
    ("多文件编码 / Agent 工作流", "Claude Opus 4.8", "最佳意图理解 + SWE-bench Pro 冠军"),
    ("Agent 终端操作 / DevOps", "GPT-5.5", "82.7% Terminal-Bench 冠军"),
    ("性价比首选 / 全科均衡", "Gemini 3.1 Pro", "80.6% SWE 得分 + 仅 $2/$12 每百万"),
    ("最低成本 (前沿品质)", "DeepSeek V4-Flash", "$0.07 / 百万输出 token，仅需对手 1/350"),
    ("自部署 / MIT 许可 / 隐私优先", "DeepSeek V4-Pro / GLM-5", "可下载权重，无数据泄露风险"),
]
for j, (scene, model, why) in enumerate(picks):
    y = Inches(5.82) + Inches(j * 0.2)
    add_text_box(slide, Inches(1.0), y, Inches(3.5), Inches(0.18), scene, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, Inches(4.6), y, Inches(3.2), Inches(0.18), f"→ {model}", font_size=11, color=ACCENT4, bold=True)
    add_text_box(slide, Inches(7.6), y, Inches(4.5), Inches(0.18), why, font_size=10, color=DIM_GRAY)

page_number(slide, 7)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — AI 风口行业 & 成功案例
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "05  AI 风口行业 & 成功案例 (一)", "金融 · 医疗 · 教育 — AI 落地进入实干期")

card(slide, Inches(0.8), Inches(1.6), Inches(3.7), Inches(2.6),
     "🏦 金融: 招商银行「AI First」",
     ["模式: 「人+Agent」协同，AI 辅助非替代",
      "成效: 2025 年实现 1556 万小时人工替代",
      "特点: 保持服务温度与合规性", "入选 2026 全国智能体十大典型案例",
      "平安保险: 7 万+智能体协同，理赔效率 ↑80%",
      "工商银行: AI 数字员工承担超 5.5 万人年工作"],
     accent_color=ACCENT)

card(slide, Inches(4.8), Inches(1.6), Inches(3.7), Inches(2.6),
     "🏥 医疗: 「天医智联」国产医疗大模型",
     ["成就: CheXpert 国际评测全球第一",
      "能力: 覆盖 200+ 疾病，胸片准确率 92.7%",
      "落地: 300+ 基层医院部署，诊断准确率 ↑40%",
      "清华 AI 医生: MedQA 准确率 96%，8 家医院公测",
      "AI 药物研发: 周期从 5 年缩至 18 个月",
      "华大基因: 全基因组分析时间大幅压缩"],
     accent_color=ACCENT2)

card(slide, Inches(8.8), Inches(1.6), Inches(3.7), Inches(2.6),
     "🎓 教育: AI 教育出海「智启教育」",
     ["融资: 红杉中国领投 1.2 亿美元 B 轮",
      "市场: 东南亚 K12，2.5 亿学生，英语教师缺口超千万",
      "用户: 500 万+，印尼/越南市场前三",
      "社会价值: 60% 用户来自农村，月费仅 10 元",
      "K12 自适应: 数学成绩平均提升 22 分",
      "高校 AI 批改: 教师批改效率提升 5×"],
     accent_color=ACCENT3)

# Key takeaway
add_rect(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.3), fill_color=BG_SECTION)
add_text_box(slide, Inches(1.0), Inches(4.6), Inches(11), Inches(0.35),
             "📌 更多风口赛道", font_size=16, color=WHITE, bold=True)

more_tracks = [
    ("🤖 具身智能/机器人", "工匠行智能签下上亿订单；宇树科技 IPO 过会，Q1 营收 ↑68.5%；特斯拉人形机器人 V3 即将发布；比亚迪确认开发人形机器人"),
    ("🐾 情感 AI / 宠物经济", "Traini 获 NVIDIA+OpenAI 高管投资 750 万美元；宠物情绪识别准确率 75-96%；美国市场 1520 亿美元"),
    ("📸 AI 消费硬件", "像素跃动 AI 写真相机，利润率 50-60%，10 人团队 5 个月从创意到产品"),
    ("🏢 政企 AI To B", "零一万物订单突破 15 亿元，2027 港股 IPO；海底捞 AI 打通全业务；元气森林 AI 驱动口味测试"),
]
for j, (track, desc) in enumerate(more_tracks):
    y = Inches(4.95) + Inches(j * 0.42)
    add_text_box(slide, Inches(1.0), y, Inches(3.0), Inches(0.2), track, font_size=12, color=ACCENT2, bold=True)
    add_text_box(slide, Inches(4.1), y, Inches(8), Inches(0.2), desc, font_size=11, color=LIGHT_GRAY)

page_number(slide, 8)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — AI 行业趋势 & 未来展望
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
section_header(slide, "06  行业趋势 & 未来展望", "2026 关键判断与战略建议")

# 5 key trends
trends = [
    ("趋势一", "从生成式 AI → 智能体 AI", "Agent 成数字生活中心，可能替代 App 生态",
     ACCENT),
    ("趋势二", "多智能体协同进入成熟期", "单一模型无法满足复杂场景，组织级 Agent 协作成标配",
     ACCENT2),
    ("趋势三", "竞争壁垒 = 行业 Know-how 数字化", "通用大模型能力收敛，行业经验数字化才是护城河",
     ACCENT3),
    ("趋势四", "中国模型性价比碾压", "DeepSeek/GLM/Qwen 开源 MIT 许可，价格仅为西方模型 1/50~1/350",
     ACCENT4),
    ("趋势五", "AI 从提效工具 → 商业模式重构", "AI 重新定义产品、服务与商业模式，情感价值成新变现点",
     RGBColor(0xE0, 0x7B, 0xFF)),
]

for i, (label, title, desc, clr) in enumerate(trends):
    y = Inches(1.6) + Inches(i * 1.05)
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.88), fill_color=BG_SECTION)
    add_rect(slide, Inches(0.8), y, Inches(0.06), Inches(0.88), fill_color=clr)
    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(1.2), Inches(0.25),
                 label, font_size=11, color=clr, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.05), Inches(9.5), Inches(0.3),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.42), Inches(9.5), Inches(0.35),
                 desc, font_size=13, color=LIGHT_GRAY)

# Strategic advice
add_rect(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.85), fill_color=BG_SECTION)
add_text_box(slide, Inches(1.0), Inches(6.15), Inches(11), Inches(0.65),
             '💡 战略建议: 2026 年是 AI 落地窗口期 — 企业应尽快将 AI Agent 嵌入核心业务流程；开发者应 & 关注 MCP/A2A 标准生态 & 优先选择性价比高的模型组合（Gemini/DeepSeek 互补）；投资者应关注垂直场景 + 软硬结合终端赛道。',
             font_size=13, color=ACCENT4, bold=True)

page_number(slide, 9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — 总结 & Thank You
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=ACCENT)
add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), fill_color=ACCENT2)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.8),
             "THANK YOU", font_size=52, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.5),
             "2026 年 6 月 AI 行业资讯报告", font_size=22, color=LIGHT_GRAY)

summary_items = [
    "✅ AI Agent 时代全面到来 — Agentic 流量已超人类",
    "✅ MCP + A2A 双协议互补 — AI 世界的 TCP/IP 标准确立",
    "✅ Claude Opus 4.8 登顶 Agent 赛道，GPT-5.5 领跑数学，Gemini 3.1 性价比最优",
    "✅ 中国开源模型崛起 — 性价比 50×+ 优势颠覆定价格局",
    "✅ 金融/医疗/教育/机器人/情感AI — 六大风口赛道已跑出商业闭环",
]
add_bullet_points(slide, Inches(1.5), Inches(3.5), Inches(9), Inches(2.2),
                  summary_items, font_size=14, color=LIGHT_GRAY)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
             "数据来源: Computex 2026 · Microsoft Build · Forbes AI 50 · Artificial Analysis · LiveBench · 阿里云案例集 · 人民日报",
             font_size=10, color=DIM_GRAY)

page_number(slide, 10)

# ── Save ───────────────────────────────────────────────────────
desktop = os.path.join(os.environ["USERPROFILE"], "Desktop")
output_path = os.path.join(desktop, "AI行业最新资讯报告_2026年6月.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
